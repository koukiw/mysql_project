import pandas as pd
from pdfminer.high_level import extract_text
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
import subprocess
import os
import xlrd


def extract_text_from_file(filepath):
    try:
        filepath_format = os.path.splitext(filepath)[1] [1:]
        text = ""
        if filepath_format =="pdf":
            text = pdf2text(filepath)
        elif filepath_format =="csv":
            text = csv2text(filepath)
        elif filepath_format =="xlsx":
            text = xlsx2text(filepath)
        elif filepath_format =="xls":
            text = xls2text(filepath)
        elif filepath_format =="docx":
            text = docx2text(filepath)
        elif filepath_format =="doc":
            text = doc2text(filepath)
        elif filepath_format =="pptx":
            text = pptx2text(filepath)
        return text
    
    except Exception as e:
        print("get_textにてerror発生")
        return -1
    

# PDFからテキストを抜き出す
def pdf2text(filepath):
    try:
        # テキストの抜き出し
        text = extract_text(filepath)
        # テキストの加工
        text = text.replace("\n","").replace("\r","").replace("\t","").strip()
        # print(text[:50])
        return text
            
    except Exception as e:
        print("pdf2textにてerror発生")
        return -1

# CSVからテキストを抜き出してjson形式で情報を整理（indexごとにjson作成）
def csv2text(filepath):
    try:
        pd_dic = pd.read_csv(filepath, sep=",")
        text = {}
        for index, dic in pd_dic.to_dict(orient="index").items():
            text["{}".format(index)] = dic
        # テキストの加工
        text = str(text)
        text = text.replace("\n","").replace("\r","").replace("\t","").strip()
        # print(text)
        return text

    except Exception as e:
        print("csv2textにてerror発生")
        return -1

#1セルごとに文字列抽出して、単語をリストに格納していく
def xlsx2text(filepath):
    try:
        workbook = load_workbook(filepath)
        sheet_names= workbook.sheetnames
        strings={}
        for index in range(len(sheet_names)):
            sheet = workbook[sheet_names[index]]

            string = []

            for row in sheet.iter_rows(values_only=True):
                for cell in row:
                    if isinstance(cell, str):
                        string.append(cell)
            strings[sheet_names[index]] = string
            

        return str(strings)
    
    except Exception as e:
        print("excel2textにてerror発生")
        return -1

#古い形式のexcelファイル（.xls）からテキスト抽出
def xls2text(filepath):
    workbook = xlrd.open_workbook(filepath)
    sheet = workbook.sheet_by_index(0) #シート情報取得
    strings= {}
    for index in range(workbook.nsheets):
        sheet = workbook.sheet_by_index(index)  
        string = []
        for row in range(sheet.nrows):
            for col in range(sheet.ncols):
                cell_value = sheet.cell_value(row, col)
                if cell_value and isinstance(cell_value, str):
                    string.append(cell_value)
        strings[sheet.name] = string
    return str(strings)


#  .docxからテキストを抜き出す
def docx2text(filepath):
    try:
        text = ""
        document = Document(filepath)
        for i, p in enumerate(document.paragraphs):
            text += p.text

        # 表中のテキスト抽出
        for i, t in enumerate(document.tables):
            for j, r in enumerate(t.rows):
                for k, c in enumerate(r.cells):
                    text += c.text
        # テキストの加工
        text = text.replace("\n","").replace("\r","").replace("\t","").strip()
        # print(text)
        return text
        
    except Exception as e:
        print("word2textにてerror発生")
        return -1

#  .docからテキストを抜き出す
def doc2text(filepath_path):
    try:
        command = f"antiword  {filepath_path}"
        # print(command)
        output = subprocess.check_output(command, shell=True)
        text = output.decode('utf-8')
        text = text.replace("\n","").replace("\r","").replace("\t","").replace(" ","").replace("　","").strip()
        # print(text)
        return text
    except subprocess.CalledProcessError:
        print("doc2textにてerror発生")
        return -1

def pptx2text(filepath):
    try:
        strings = []
        prs = Presentation(filepath)
        for i, sld in enumerate(prs.slides, start=1):
            for shp in sld.shapes:
                if shp.has_text_frame:
                    strings.append(shp.text)
                #テーブル（表）からも文字列抽出
                if shp.has_table:
                    tables =[]
                    for row in shp.table.rows:
                        values = [ cell.text for cell in row.cells ]
                        tables.append(values)
                    strings.append(tables)
        return str(strings)
    except Exception as e:
        print("pptx2textにてerror発生")
        return -1

# メイン処理
if __name__ == "__main__":
    # print("pdfminer.six ver.{}".format(pdfminer.__version__))
    results = pdf2text()
    print(results)

