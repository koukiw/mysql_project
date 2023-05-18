from pptx import Presentation


# create_by_chatGPT
# def extract_strings_from_powerpoint(file_path):
#     presentation = Presentation(file_path)
#     for slide in presentation.slides:
#         print("slide",slide)
#     strings = []

#     for slide in presentation.slides:
#         for shape in slide.shapes:
#             print("shape",shape)
#             if shape.has_text_frame:
#                 for paragraph in shape.text_frame.paragraphs:
#                     for run in paragraph.runs:
#                         strings.append(run.text)

#     return strings


def pptx2text(file_path):
    strings = []
    prs = Presentation(file_path)
    for i, sld in enumerate(prs.slides, start=1):
        # print(f'-- {i} --')
        for shp in sld.shapes:
            if shp.has_text_frame:
                # print("shp",shp)
                # print("shp.text",shp.text)
                strings.append(shp.text)
            if shp.has_table:
                # print("table発見")
                tables =[]
                for row in shp.table.rows:
                    values = [ cell.text for cell in row.cells ]
                    tables.append(values)
                strings.append(tables)
    return strings

# PowerPointファイルのパスを指定して実行
file_path = './DBプラットフォーム検討.pptx'
# result = extract_strings_from_powerpoint(file_path)
# print(result)
result = pptx2text(file_path)
print(result)
